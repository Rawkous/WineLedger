"""
Generate a 9-slide PowerPoint pitch deck for WineLedger (Phases A–C).

Usage (from repo root):
  python3 -m venv .venv-pptx
  . .venv-pptx/bin/activate
  python -m pip install -r doc/requirements-pptx.txt
  python doc/build_pitch_deck_pptx.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT_PPTX = Path(__file__).resolve().parent / "wineledger-phases-a-c-deck.pptx"


def _set_text_run(run, *, size_pt: int | None = None, bold: bool = False, rgb: tuple[int, int, int] | None = None):
    if size_pt is not None:
        run.font.size = int(size_pt * 12700)  # EMU per point
    run.font.bold = bool(bold)
    if rgb is not None:
        run.font.color.rgb = RGBColor(*rgb)


def add_title(prs: Presentation, title: str, subtitle: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # title slide
    slide.shapes.title.text = title
    if subtitle is not None:
        slide.placeholders[1].text = subtitle


def add_bullets(prs: Presentation, title: str, bullets: list[str], *, note: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = 22 * 12700
    if note:
        tx = slide.shapes.add_textbox(left=prs.slide_width * 0.07, top=prs.slide_height * 0.83, width=prs.slide_width * 0.86, height=prs.slide_height * 0.12)
        tf = tx.text_frame
        tf.text = note
        tf.paragraphs[0].font.size = 14 * 12700
        tf.paragraphs[0].font.italic = True
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT


def build() -> None:
    prs = Presentation()

    add_title(
        prs,
        "WineLedger",
        "A living digital twin of the wine supply chain — Phases A–C (Cardano route)",
    )

    add_bullets(
        prs,
        "The problem",
        [
            "Supply-chain stories are hard to verify after the fact",
            "Data is siloed; partners/auditors lack a shared, timestamped view",
            "Most systems show spreadsheets — not a story people can feel",
        ],
        note="WineLedger turns provenance into an append-only event story + real-time generative visuals.",
    )

    add_bullets(
        prs,
        "What we built (Phase A — complete)",
        [
            "FastAPI backend with a persistent, hash-linked event chain (JSON on disk)",
            "WebSockets stream new blocks live to clients",
            "Vite + p5.js generative canvas + education page (harvest → retail)",
            "Optional geo enrichment via pluggable cache (SQLite)",
        ],
    )

    add_bullets(
        prs,
        "Phase B — production + Cardano anchoring",
        [
            "Deploy WineLedger in a private cloud environment (AWS VPC + TLS + ops hardening)",
            "Durable exports/snapshots in object storage (S3-class) with content hashes",
            "Cardano tx builder + submitter (isolated signer) posts verifiable anchors",
            "Public verification via explorers/wallet tooling (no proprietary verifier required)",
        ],
        note="Goal: integrity is independently checkable, not just “inside our server.”",
    )

    add_bullets(
        prs,
        "Phase C — culture layer (media + NFTs + installations)",
        [
            "A.I.-assisted imagery and campaign media derived from real milestones",
            "Cardano-native NFT drops using CIP-25 metadata + time-locked policies",
            "Long-form generative video for tasting rooms/retail screens",
            "Everything still points back to verifiable snapshots and the core ledger",
        ],
    )

    add_bullets(
        prs,
        "Cardano strategy (two routes)",
        [
            "Route A: NFTs without smart contracts (native scripts + CIP-25 / label 721)",
            "Route B: long-lived identity using CIP-68 (reference + updatable state)",
            "Hydra heads optionally handle high-frequency updates; settle to L1 at milestones",
            "Standards-first approach for interoperability and lower operational risk",
        ],
    )

    add_bullets(
        prs,
        "How the system fits together",
        [
            "WineLedger API remains the source of truth for detailed events",
            "Exports/snapshots are hashed and stored durably",
            "Cardano anchors provide timestamped public commitments to those hashes",
            "Visual layer turns each event into motion/color/rhythm in real time",
        ],
    )

    add_bullets(
        prs,
        "Demo flow (what stakeholders see)",
        [
            "Run a simulated journey: harvest → fermentation → aging → bottling → transport → retail",
            "Watch blocks stream live and the canvas evolve with each event",
            "Show an exported snapshot + its hash",
            "Point to the Cardano anchor proving when that snapshot existed",
        ],
    )

    add_bullets(
        prs,
        "Next steps",
        [
            "Finalize Cardano network + operational custody model (testnet → mainnet)",
            "Stand up tx-builder sidecar + monitoring; begin milestone anchoring",
            "Prepare first CIP-25 collection drop tied to a real release/campaign",
            "Pilot with one partner route; expand to broader participants",
        ],
        note="Deliverable: a working product + a verifiable public provenance layer + a compelling visual narrative.",
    )

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)
    print(f"Wrote: {OUT_PPTX}")


if __name__ == "__main__":
    build()

